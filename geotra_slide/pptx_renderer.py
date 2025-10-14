"""Utilities to render :class:`SlideDocument` objects into PPTX files."""

from __future__ import annotations

import io
import subprocess
import tempfile
from pathlib import Path
from typing import Dict, Iterable, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import SlidePlaceholder

from .slide_library import SlideLibrary
from .slide_models import SlideDocument, SlidePlaceholderContent


class SlideDeckRenderer:
    """Render slide documents into PPTX binaries (and optional previews)."""

    def __init__(self, slide_library: SlideLibrary) -> None:
        self.slide_library = slide_library
        self.master_template_path = self.slide_library.master_template_path()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------
    def render_document(self, document: SlideDocument) -> io.BytesIO:
        """Return a PPTX stream that represents ``document``."""

        presentation = Presentation(self.master_template_path)
        _clear_existing_slides(presentation)

        for slide_page in document.slides:
            asset = self.slide_library.get_asset(slide_page.asset_id)
            source_path = self.slide_library.asset_file_path(asset.asset_id)
            source_prs = Presentation(source_path)
            template_slide = self._copy_slide(source_prs, presentation, 0)
            self._write_placeholders(template_slide, slide_page)

        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        return buffer

    def render_preview_image(
        self,
        document: SlideDocument,
        *,
        slide_index: int = 0,
        pptx_bytes: Optional[bytes] = None,
    ) -> Optional[bytes]:
        """Generate a PNG preview for ``document`` if LibreOffice is available."""

        try:
            payload = pptx_bytes or self.render_document(document).getvalue()
            with tempfile.TemporaryDirectory() as tmpdir:
                pptx_path = Path(tmpdir) / "preview.pptx"
                pptx_path.write_bytes(payload)

                soffice_path = _locate_soffice()
                if soffice_path is None:
                    return None

                cmd = [
                    soffice_path,
                    "--headless",
                    "--convert-to",
                    "png",
                    "--outdir",
                    tmpdir,
                    str(pptx_path),
                ]
                subprocess.run(
                    cmd,
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )

                png_files = sorted(Path(tmpdir).glob("*.png"))
                if not png_files:
                    return None

                index = max(0, min(slide_index, len(png_files) - 1))
                return png_files[index].read_bytes()
        except Exception:
            return None

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------
    def _copy_slide(
        self, source_prs: Presentation, destination_prs: Presentation, slide_index: int
    ):
        source_slide = source_prs.slides[slide_index]
        layout_name = source_slide.slide_layout.name

        try:
            layout = destination_prs.slide_layouts.get_by_name(layout_name)
        except KeyError:
            layout = destination_prs.slide_layouts[1]

        new_slide = destination_prs.slides.add_slide(layout)
        self._clone_non_placeholder_shapes(source_slide.shapes, new_slide.shapes)

        required_indices = {
            shape.placeholder_format.idx
            for shape in source_slide.shapes
            if shape.is_placeholder
        }
        present_indices = {
            shape.placeholder_format.idx
            for shape in new_slide.shapes
            if shape.is_placeholder
        }
        missing = required_indices - present_indices
        if missing:
            raise RuntimeError(
                "Master template is missing placeholders: " f"{sorted(missing)}"
            )

        return new_slide

    def _clone_non_placeholder_shapes(
        self, source_shapes: Iterable[BaseShape], destination_shapes
    ) -> None:
        for shape in source_shapes:
            if shape.is_placeholder:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                destination_shapes.add_picture(
                    io.BytesIO(shape.image.blob),
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )
            elif shape.has_text_frame:
                new_shape = destination_shapes.add_textbox(
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height,
                )
                if shape.text_frame.text:
                    new_shape.text_frame.text = shape.text_frame.text

    def _write_placeholders(self, slide, slide_page) -> None:
        placeholders_by_idx: Dict[int, SlidePlaceholder] = {}
        for shape in slide.shapes:
            if shape.is_placeholder:
                placeholders_by_idx[shape.placeholder_format.idx] = shape

        asset = self.slide_library.get_asset(slide_page.asset_id)
        placeholder_specs = {spec.name: spec for spec in asset.placeholders}
        content_map = {item.name: item for item in slide_page.placeholders}

        for name, spec in placeholder_specs.items():
            placeholder = placeholders_by_idx.get(spec.idx)
            if placeholder is None:
                continue
            text = content_map.get(name, SlidePlaceholderContent(name, spec.description, spec.edit_policy)).text
            if placeholder.has_text_frame:
                text_frame = placeholder.text_frame
                text_frame.clear()
                text_frame.text = text

        if slide_page.title:
            title_shape = next(
                (
                    shape
                    for shape in slide.shapes
                    if shape.is_placeholder
                    and shape.placeholder_format.type == 1  # TITLE
                ),
                None,
            )
            if title_shape and title_shape.has_text_frame:
                title_shape.text_frame.text = slide_page.title


# ----------------------------------------------------------------------
# Helper functions
# ----------------------------------------------------------------------

def _clear_existing_slides(presentation: Presentation) -> None:
    for idx in range(len(presentation.slides) - 1, -1, -1):
        slide_id = presentation.slides._sldIdLst[idx].rId
        presentation.part.drop_rel(slide_id)
        del presentation.slides._sldIdLst[idx]


def _locate_soffice() -> Optional[str]:
    candidates = [
        "soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
    ]
    for candidate in candidates:
        try:
            subprocess.run(
                [candidate, "--version"],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            return candidate
        except Exception:
            continue
    return None
