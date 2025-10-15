import io
from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation

from geotra_slide.pptx_renderer import SlideDeckRenderer
from geotra_slide.slide_library import SlideLibrary
from geotra_slide.slide_models import (
    SlideDocument,
    SlidePage,
    SlidePlaceholderContent,
)


def _build_document(library: SlideLibrary) -> SlideDocument:
    assets = sorted(library.list_assets(), key=lambda item: item.asset_id)
    asset = next(
        asset
        for asset in assets
        if any(ph.edit_policy.lower() == "generate" for ph in asset.placeholders)
    )
    placeholders = [
        SlidePlaceholderContent(name=spec.name, text=f"テスト:{spec.name}", policy=spec.edit_policy)
        for spec in asset.placeholders
    ]
    slide = SlidePage(
        slide_id="slide_01",
        page_number=1,
        asset_id=asset.asset_id,
        asset_file=asset.file_name,
        title="テストタイトル",
        placeholders=placeholders,
    )
    return SlideDocument(slides=[slide])


def test_render_document_populates_placeholders():
    library = SlideLibrary(Path("assets"))
    renderer = SlideDeckRenderer(library)
    document = _build_document(library)

    pptx_stream = renderer.render_document(document)
    prs = Presentation(io.BytesIO(pptx_stream.getvalue()))

    texts = [
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert any(text.startswith("テスト:") for text in texts)


def test_render_preview_image_returns_none_without_soffice(monkeypatch):
    library = SlideLibrary(Path("assets"))
    renderer = SlideDeckRenderer(library)
    document = _build_document(library)

    pptx_stream = renderer.render_document(document)
    monkeypatch.setattr("geotra_slide.pptx_renderer._locate_soffice", lambda: None)
    preview = renderer.render_preview_image(document, pptx_bytes=pptx_stream.getvalue())
    assert preview is None
