"""End-to-end I/O tests for the slide generation pipeline."""

from __future__ import annotations

from pathlib import Path

import pytest

pytest.importorskip("pptx")
from pptx import Presentation

from geotra_slide.pptx_renderer import SlideDeckRenderer
from geotra_slide.slide_generation import (
    GenerationContext,
    PlanningContext,
    SlideContentGenerator,
    SlideOutlineGenerator,
    SlideStructurePlanner,
)
from geotra_slide.slide_library import SlideLibrary
from tests.llm_stubs import MultiStageStubLLM


@pytest.fixture(scope="module")
def slide_library() -> SlideLibrary:
    return SlideLibrary(Path("assets"))


def test_pipeline_generates_multi_slide_pptx(tmp_path, slide_library):
    """Ensure dummy inputs flowing through the pipeline produce PPTX output."""

    outline_payload = {
        "slides": [
            {
                "slide_id": "slide_01",
                "page_number": 1,
                "asset_id": "cover_regular_001",
                "title": "提案の概要",
                "notes": "冒頭で目的を整理",
            },
            {
                "slide_id": "slide_02",
                "page_number": 2,
                "asset_id": "schedule_001",
                "title": "実行ロードマップ",
                "notes": "四半期計画を提示",
            },
        ]
    }
    placeholder_payloads = [
        {
            "placeholders": [
                {
                    "placeholder_name": "テキスト プレースホルダー 3",
                    "text": "ACME社向け成長戦略の骨子を提示",
                    "references": ["internal_report.md"],
                }
            ],
            "slide_summary": "提案背景",
            "citations": ["internal_report.md"],
        },
        {
            "placeholders": [
                {
                    "placeholder_name": "テキスト プレースホルダー 3",
                    "text": "次の四半期までのアクションアイテム",
                    "references": ["https://example.com/roadmap"],
                }
            ],
            "slide_summary": "計画概要",
            "citations": ["https://example.com/roadmap"],
        },
    ]

    stub_llm = MultiStageStubLLM(
        outline_payload=outline_payload,
        placeholder_payloads=placeholder_payloads,
        structure_text="提案概要とロードマップの二部構成",
        web_search_text="外部リサーチ: 最新市場レポート",
    )

    planner = SlideStructurePlanner(stub_llm)
    planning_context = PlanningContext(
        conversation_history="ユーザー: ACME社向け提案資料を作りたい",
        goal="営業提案の骨子をまとめる",
        target_company="ACME社",
    )
    structure_text = planner.build_structure(planning_context)
    assert "二部構成" in structure_text

    outline_context = GenerationContext(
        user_request="ACME社に提案する資料を準備したい",
        target_company="ACME社",
        additional_notes="冒頭で狙いを整理する",
    )
    document = SlideOutlineGenerator(slide_library, llm_client=stub_llm).generate_outline(
        slide_structure=structure_text,
        context=outline_context,
    )
    assert len(document.slides) == 2

    generation_context = GenerationContext(
        user_request="ACME社向け提案資料を作成したい",
        target_company="ACME社",
        additional_notes="重要メッセージは簡潔に",
        perform_web_search=True,
    )
    populated_document = SlideContentGenerator(
        slide_library,
        llm_client=stub_llm,
        internal_document_path=Path("data/internal_report.md"),
    ).generate_for_document(document, context=generation_context)

    assert populated_document.metadata["references"] == [
        "https://example.com/roadmap",
        "internal_report.md",
    ]
    assert stub_llm.web_search_requests, "Web search should be invoked when enabled"

    renderer = SlideDeckRenderer(slide_library)
    pptx_stream = renderer.render_document(populated_document)
    output_path = tmp_path / "pipeline_output.pptx"
    output_path.write_bytes(pptx_stream.getvalue())

    prs = Presentation(output_path)
    assert len(prs.slides) == 2
    slide_texts = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert any("成長戦略" in text for text in slide_texts)
    assert any("アクションアイテム" in text for text in slide_texts)
