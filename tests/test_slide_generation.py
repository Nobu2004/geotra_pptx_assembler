import io
from pathlib import Path

import pytest

pytest.importorskip("pptx")

from pptx import Presentation

from geotra_slide.pptx_renderer import SlideDeckRenderer
from geotra_slide.slide_document import SlideDocumentStore
from geotra_slide.slide_generation import (
    GenerationContext,
    PlanningContext,
    SlideContentGenerator,
    SlideOutlineGenerator,
    SlideStructurePlanner,
)
from geotra_slide.slide_library import SlideLibrary

from tests.llm_stubs import MultiStageStubLLM


@pytest.fixture(scope="module")
def slide_library():
    return SlideLibrary(Path("assets"))


def test_multi_slide_generation_workflow(slide_library):
    outline_payload = {
        "slides": [
            {
                "slide_id": "slide_01",
                "page_number": 1,
                "asset_id": "cover_regular_001",
                "title": "定例報告の目的",
                "notes": "冒頭で期待値を整理",
            },
            {
                "slide_id": "slide_02",
                "page_number": 2,
                "asset_id": "schedule_001",
                "title": "主要スケジュール",
            },
        ]
    }
    placeholder_payloads = [
        {
            "placeholders": [
                {
                    "placeholder_name": "テキスト プレースホルダー 3",
                    "text": "JKAの進捗を網羅的に提示",
                    "references": ["internal_report.md"],
                }
            ],
            "slide_summary": "JKA向け進捗概要",
            "citations": ["internal_report.md"],
        },
        {
            "placeholders": [
                {
                    "placeholder_name": "テキスト プレースホルダー 3",
                    "text": "次回会議で確認するトピック",
                    "references": ["https://example.com/schedule"],
                }
            ],
            "slide_summary": "打ち合わせ計画",
            "citations": ["https://example.com/schedule"],
        },
    ]

    stub_llm = MultiStageStubLLM(
        outline_payload=outline_payload,
        placeholder_payloads=placeholder_payloads,
        structure_text="表紙とスケジュールを含む2枚構成で提案",
    )

    planner = SlideStructurePlanner(stub_llm)
    planning_context = PlanningContext(
        conversation_history="ユーザー: 進捗報告の構成を相談したい",
        goal="定例ミーティングで使う資料の骨子を作る",
        target_company="JKA",
    )
    structure_text = planner.build_structure(planning_context)
    assert "2枚構成" in structure_text

    outline_context = GenerationContext(
        user_request="定例ミーティング向け資料を作りたい",
        target_company="JKA",
        additional_notes="冒頭で目的を揃える",
    )
    outline_generator = SlideOutlineGenerator(slide_library, llm_client=stub_llm)
    document = outline_generator.generate_outline(
        slide_structure=structure_text,
        context=outline_context,
    )
    assert len(document.slides) == 2
    assert document.metadata["slide_structure"] == structure_text

    content_generator = SlideContentGenerator(
        slide_library,
        llm_client=stub_llm,
        internal_document_path=Path("data/internal_report.md"),
    )
    generation_context = GenerationContext(
        user_request="定例ミーティング向け資料を作りたい",
        target_company="JKA",
        external_research=None,
        additional_notes="課題感を明確にする",
        perform_web_search=True,
    )
    updated_document = content_generator.generate_for_document(
        document,
        context=generation_context,
    )

    assert updated_document.metadata["references"] == [
        "https://example.com/schedule",
        "internal_report.md",
    ]

    slide1 = updated_document.get_slide("slide_01")
    assert slide1 is not None
    assert slide1.notes["summary"] == "JKA向け進捗概要"
    placeholder_map = {item.name: item for item in slide1.placeholders}
    assert placeholder_map["テキスト プレースホルダー 2"].text == "JKA"
    assert placeholder_map["テキスト プレースホルダー 3"].text.startswith("JKAの進捗")

    slide2 = updated_document.get_slide("slide_02")
    assert slide2 is not None
    assert slide2.notes["summary"] == "打ち合わせ計画"
    placeholder_map2 = {item.name: item for item in slide2.placeholders}
    assert placeholder_map2["テキスト プレースホルダー 3"].references == [
        "https://example.com/schedule"
    ]


def test_slide_document_store_roundtrip(tmp_path):
    from geotra_slide.slide_models import SlideDocument, SlidePage

    slide = SlidePage(
        slide_id="slide_99",
        page_number=1,
        asset_id="dummy",
        asset_file="dummy.pptx",
        title="テスト",
    )
    document = SlideDocument(slides=[slide], metadata={"references": []})

    path = tmp_path / "slide.json"
    store = SlideDocumentStore(path)
    store.save(document)

    loaded = store.load()
    assert loaded.to_dict() == document.to_dict()


def test_placeholder_generation_enables_pptx_rendering(slide_library):
    outline_payload = {
        "slides": [
            {
                "slide_id": "slide_01",
                "page_number": 1,
                "asset_id": "cover_regular_001",
                "title": "四半期の狙い",
            }
        ]
    }
    placeholder_payloads = [
        {
            "placeholders": [
                {
                    "placeholder_name": "テキスト プレースホルダー 3",
                    "text": "四半期概況を冒頭で整理",
                    "references": ["internal_report.md"],
                }
            ],
            "slide_summary": "QBR概要",
            "citations": ["internal_report.md"],
        }
    ]

    stub_llm = MultiStageStubLLM(
        outline_payload=outline_payload,
        placeholder_payloads=placeholder_payloads,
        structure_text="表紙とハイライトで構成",
    )

    outline_context = GenerationContext(
        user_request="四半期事業レビュー資料を作成したい",
        target_company="ACME",
    )
    outline_generator = SlideOutlineGenerator(slide_library, llm_client=stub_llm)
    document = outline_generator.generate_outline(
        slide_structure="四半期レビューの章立て", context=outline_context
    )

    content_generator = SlideContentGenerator(
        slide_library,
        llm_client=stub_llm,
        internal_document_path=Path("data/internal_report.md"),
    )
    generation_context = GenerationContext(
        user_request="四半期事業レビュー資料を作成したい",
        target_company="ACME",
        perform_web_search=False,
    )
    populated_document = content_generator.generate_for_document(
        document, context=generation_context
    )

    renderer = SlideDeckRenderer(slide_library)
    pptx_stream = renderer.render_document(populated_document)
    prs = Presentation(io.BytesIO(pptx_stream.getvalue()))
    assert len(prs.slides) == 1
    texts = [
        shape.text_frame.text
        for shape in prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert any("四半期概況" in text for text in texts)
