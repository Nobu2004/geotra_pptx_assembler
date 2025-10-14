# src/core/pptx_utils.py

import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import subprocess
import platform
import tempfile
import shutil

# --- 画像変換に関する設定 ---
# LibreOfficeに含まれる 'soffice' コマンドを主要な変換ツールとしてチェック
CAN_CONVERT_TO_IMAGE = platform.system() == "Windows" or shutil.which("soffice") is not None

def get_layouts_from_master(pptx_path: Path) -> list[dict]:
    """
    マスターテンプレートから全レイアウトの情報を抽出する。
    AIが「generated」タイプのスライドを生成する際のカタログとなる。
    """
    try:
        prs = Presentation(pptx_path)
        layouts_info = []
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = [
                {
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                    "idx": ph.placeholder_format.idx,
                }
                for ph in layout.placeholders
            ]
            layouts_info.append({
                "layout_index": i,
                "layout_name": layout.name,
                "placeholders": sorted(placeholders, key=lambda p: p["idx"]),
            })
        return layouts_info
    except Exception as e:
        print(f"エラー: {pptx_path} の読み込みに失敗しました。詳細: {e}")
        return []

def get_placeholders_from_slide(pptx_path: Path, slide_index: int = 0) -> list[dict]:
    """
    指定されたスライドからプレースホルダーの情報を抽出する。
    サンプルスライド登録時に利用。
    """
    try:
        prs = Presentation(pptx_path)
        if not prs.slides or len(prs.slides) <= slide_index:
            return []
        slide = prs.slides[slide_index]
        
        sorted_shapes = sorted(
            [shp for shp in slide.shapes if shp.is_placeholder],
            key=lambda shp: shp.placeholder_format.idx
        )
        return [{"name": shp.name, "idx": shp.placeholder_format.idx} for shp in sorted_shapes]
    except Exception as e:
        print(f"エラー: {pptx_path} の読み込みに失敗しました。詳細: {e}")
        return []

def create_annotated_preview(original_path: Path, output_image_path: Path) -> bool:
    """
    スライドのプレースホルダーにその名前を書き込み、画像として出力する。
    """
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_pptx_path = Path(temp_dir) / "annotated.pptx"
        
        try:
            prs = Presentation(original_path)
            if not prs.slides:
                return False

            slide = prs.slides[0]
            for shape in slide.shapes:
                if shape.is_placeholder and shape.has_text_frame:
                    text_frame = shape.text_frame
                    text_frame.clear()
                    p = text_frame.paragraphs[0]
                    run = p.add_run()
                    run.text = f"[{shape.name}]\n(idx: {shape.placeholder_format.idx})"
            
            prs.save(temp_pptx_path)
            
            return _convert_pptx_to_image(temp_pptx_path, output_image_path)

        except Exception as e:
            print(f"注釈付きプレビューの生成中にエラーが発生しました: {e}")
            return False

def _convert_pptx_to_image(pptx_path: Path, output_image_path: Path) -> bool:
    """
    PPTXの最初のスライドをPNG画像に変換する内部関数。
    """
    soffice_cmd = shutil.which("soffice")

    if soffice_cmd:
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                temp_output_dir = Path(temp_dir)
                subprocess.run(
                    [
                        soffice_cmd,
                        "--headless",
                        "--convert-to", "png",
                        "--outdir", str(temp_output_dir),
                        str(pptx_path)
                    ],
                    check=True, capture_output=True, timeout=30
                )
                
                generated_file = temp_output_dir / (pptx_path.stem + ".png")
                
                if generated_file.exists():
                    generated_file.rename(output_image_path)
                    return True
                else:
                    print(f"エラー: sofficeは実行されましたが、期待された出力ファイルが見つかりませんでした。")
                    return False

        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired) as e:
            print(f"sofficeでの画像変換に失敗しました: {e}")
            if isinstance(e, subprocess.CalledProcessError):
                print(f"Stderr: {e.stderr.decode()}")
            return False
            
    elif platform.system() == "Windows":
        print("Windowsでの画像変換は未実装です。")
        return False
        
    else:
        print("警告: 画像変換機能が利用できません (sofficeコマンドが見つかりません)。")
        print("解決策: `brew install --cask libreoffice` を実行してLibreOfficeをインストールしてください。")
        return False
