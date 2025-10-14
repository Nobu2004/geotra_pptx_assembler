# src/core/renderer.py

import json
import io
import os
import glob
import tempfile
import subprocess
from pathlib import Path
from typing import List, Dict

from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

from . import schemas

ROOT_DIR = Path(__file__).parent.parent.parent
ASSETS_DIR = ROOT_DIR / "assets"
SLIDE_LIBRARY_DIR = ASSETS_DIR / "slide_library"
TEMPLATE_DIR = ASSETS_DIR / "templates"

class PPTXRenderer:
    def __init__(self):
        master_manifest_path = TEMPLATE_DIR / "master_manifest.json"
        self.slide_lib_manifest_path = SLIDE_LIBRARY_DIR / "slide_library_manifest.json"

        with open(master_manifest_path, 'r', encoding='utf-8') as f:
            self.master_manifest = json.load(f)
        
        if not self.slide_lib_manifest_path.exists():
            raise FileNotFoundError(f"Slide library manifest not found at: {self.slide_lib_manifest_path}")

        with open(self.slide_lib_manifest_path, 'r', encoding='utf-8') as f:
            self.slide_library_manifest = json.load(f)
            self.slide_asset_map: Dict[str, dict] = {
                asset['id']: asset for asset in self.slide_library_manifest.get('slide_assets', [])
            }
        
        self.master_template_path = TEMPLATE_DIR / self.master_manifest['master_template_file']

    def render_presentation(self, blueprints: List[schemas.SlideBlueprint]) -> io.BytesIO:
        prs = Presentation(self.master_template_path)
        
        if len(prs.slides) > 0:
            for i in range(len(prs.slides) - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]

        for blueprint in blueprints:
            try:
                print(f"\n--- Blueprint '{blueprint.asset_id}' のスライドを生成中... ---")
                self._create_slide_from_blueprint(prs, blueprint)
                print(f"--- Blueprint '{blueprint.asset_id}' の生成完了 ---")
            except Exception as e:
                self._add_error_slide(prs, blueprint, e)
        
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io

    def render_slide_preview_image(self, blueprints: List[schemas.SlideBlueprint], slide_index: int) -> bytes | None:
        """
        選択スライドのプレビュー画像(PNG)を返す。
        LibreOffice(soffice)がインストールされていない場合は None を返す。
        """
        try:
            pptx_io = self.render_presentation(blueprints)
            with tempfile.TemporaryDirectory() as tmpdir:
                pptx_path = os.path.join(tmpdir, "preview.pptx")
                with open(pptx_path, "wb") as f:
                    f.write(pptx_io.getvalue())

                # soffice のパス候補
                soffice_candidates = [
                    "soffice",
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                ]
                soffice_path = next((p for p in soffice_candidates if _is_executable(p)), None)
                if not soffice_path:
                    return None

                cmd = [
                    soffice_path,
                    "--headless",
                    "--convert-to", "png",
                    "--outdir", tmpdir,
                    pptx_path,
                ]
                subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)

                png_files = sorted(glob.glob(os.path.join(tmpdir, "*.png")))
                if not png_files:
                    return None
                idx = max(0, min(slide_index, len(png_files)-1))
                with open(png_files[idx], "rb") as imgf:
                    return imgf.read()
        except Exception as e:
            print(f"[情報] スライドプレビュー画像の生成に失敗: {e}")
            return None


def _is_executable(path: str) -> bool:
    try:
        # PATH 上のコマンド名にも対応
        return subprocess.call([path, "--version"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL) == 0
    except Exception:
        return False

class PPTXRenderer(PPTXRenderer):
    # Reattach helper methods that were accidentally outdented
    def _create_slide_from_blueprint(self, prs: Presentation, blueprint: schemas.SlideBlueprint):
        asset_info = self.slide_asset_map.get(blueprint.asset_id)
        if not asset_info:
            raise ValueError(f"アセットID '{blueprint.asset_id}' がマニフェストに見つかりません。")

        source_pptx_path = SLIDE_LIBRARY_DIR / asset_info['file_name']
        source_prs = Presentation(source_pptx_path)
        
        new_slide = self._copy_slide(source_prs, prs, 0)

        manifest_placeholders_by_name: Dict[str, dict] = {
            ph['name']: ph for ph in asset_info.get('placeholders', [])
        }
        placeholders_on_slide_by_idx: Dict[int, SlidePlaceholder] = {
            shape.placeholder_format.idx: shape for shape in new_slide.shapes if shape.is_placeholder
        }
        
        print(f"コンテンツをプレースホルダーに書き込み中...")
        for content_item in blueprint.content_map:
            placeholder_name = content_item.placeholder_name
            ph_def = manifest_placeholders_by_name.get(placeholder_name)
            
            if not ph_def:
                print(f"  [警告] マニフェストに '{placeholder_name}' の定義が見つかりません。")
                continue
            
            placeholder_idx = ph_def.get('idx')
            if placeholder_idx is None:
                print(f"  [警告] マニフェストの '{placeholder_name}' に 'idx' がありません。")
                continue

            if placeholder_idx in placeholders_on_slide_by_idx:
                shape = placeholders_on_slide_by_idx[placeholder_idx]
                if shape.has_text_frame:
                    shape.text_frame.text = content_item.content
                    print(f"  ✓ idx:{placeholder_idx} ('{placeholder_name}') にコンテンツを書き込みました。")
                else:
                    print(f"  [警告] プレースホルダー (idx: {placeholder_idx}) にテキストフレームがありません。")
            else:
                # この警告が根本原因を示唆
                print(f"  [警告] 生成されたスライドにインデックス {placeholder_idx} ('{placeholder_name}') のプレースホルダーが見つかりませんでした。")

    # --- ▼▼▼ デバッグ機能強化箇所 ▼▼▼ ---
    def _copy_slide(self, prs_from: Presentation, prs_to: Presentation, slide_index: int):
        source_slide = prs_from.slides[slide_index]
        source_layout_name = source_slide.slide_layout.name
        print(f"コピー元スライドのレイアウト名: '{source_layout_name}'")

        try:
            slide_layout = prs_to.slide_layouts.get_by_name(source_layout_name)
            print(f"マスターテンプレートから対応レイアウト '{source_layout_name}' を発見しました。")
        except KeyError:
            print(f"!!! [重大な警告] レイアウト '{source_layout_name}' がマスターに見つかりません。デフォルトレイアウトを使用します。")
            slide_layout = prs_to.slide_layouts[1] 
        
        new_slide = prs_to.slides.add_slide(slide_layout)

        # --- 不整合チェックと詳細なデバッグ出力 ---
        source_ph_indices = {shape.placeholder_format.idx for shape in source_slide.shapes if shape.is_placeholder}
        new_slide_ph_indices = {shape.placeholder_format.idx for shape in new_slide.shapes if shape.is_placeholder}
        
        print(f"コピー元スライドのプレースホルダーidx: {sorted(list(source_ph_indices))}")
        print(f"生成スライドのプレースホルダーidx:   {sorted(list(new_slide_ph_indices))}")

        missing_indices = source_ph_indices - new_slide_ph_indices
        if missing_indices:
            raise RuntimeError(
                f"レイアウトの不整合を検出しました！\n"
                f"サンプルスライド '{prs_from.part.partname}' はレイアウト '{source_layout_name}' を使用しており、\n"
                f"このレイアウトはプレースホルダーidx {sorted(list(missing_indices))} を必要とします。\n"
                f"しかし、マスターテンプレート '{prs_to.part.partname}' の同名レイアウトにはこれらのプレースホルダーが存在しません。\n"
                f"マスターテンプレートを修正するか、サンプルスライドが正しいレイアウトを使用しているか確認してください。"
            )
        
        # --- プレースホルダー以外の図形をコピー ---
        for shape in source_slide.shapes:
            if not shape.is_placeholder:
                self._clone_shape(shape, new_slide.shapes)
                
        return new_slide

    def _clone_shape(self, shape, shapes_collection):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shapes_collection.add_picture(
                io.BytesIO(shape.image.blob),
                shape.left, shape.top, shape.width, shape.height
            )
        elif shape.has_text_frame:
            new_shape = shapes_collection.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            if shape.text_frame.text:
                new_shape.text_frame.text = shape.text_frame.text
            return new_shape
        return None
    # --- ▲▲▲ 修正ここまで ▲▲▲ ---

    def _add_error_slide(self, prs: Presentation, blueprint: schemas.SlideBlueprint, error: Exception):
        try:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            
            title_shape = slide.shapes.title if slide.shapes.has_title else None
            body_shape = None
            # ボディプレースホルダーをより確実に見つける
            for shape in slide.placeholders:
                # 1はBODY, 7はOBJECT(これもテキストを持てる)
                if shape.placeholder_format.type in (MSO_SHAPE_TYPE.BODY, MSO_SHAPE_TYPE.OBJECT):
                    body_shape = shape
                    break

            if title_shape:
                title_shape.text = f"エラー: スライド '{blueprint.slide_title}' の生成に失敗"
            
            if body_shape:
                body = body_shape.text_frame
                body.clear() # 既存のテキストをクリア
                body.text = "以下のエラーが発生しました:\n"
                p = body.add_paragraph()
                p.text = str(error)
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(255, 0, 0)
        except Exception as e_final:
            print(f"!!! 重大なエラー: エラースライドの生成に失敗しました: {e_final}")

