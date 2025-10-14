# tests/core/test_renderer.py

import unittest
import tempfile
from pathlib import Path
from pptx import Presentation
import json

from src.core.renderer import PPTXRenderer
from src.core import schemas

class TestPPTXRenderer(unittest.TestCase):
    def setUp(self):
        try:
            self.renderer = PPTXRenderer()
            with open(self.renderer.slide_lib_manifest_path, 'r', encoding='utf-8') as f:
                self.slide_library_manifest = json.load(f)
        except FileNotFoundError as e:
            self.fail(f"テストの前提条件エラー: マニフェストファイルが見つかりません。 {e}")
        
        self.temp_dir = tempfile.TemporaryDirectory()
        self.output_path = Path(self.temp_dir.name)

    def tearDown(self):
        self.temp_dir.cleanup()

    def test_render_single_slide_presentation(self):
        asset_id_to_test = "org_chart_001"
        asset_info = next((asset for asset in self.slide_library_manifest['slide_assets'] if asset['id'] == asset_id_to_test), None)
        self.assertIsNotNone(asset_info, f"テスト用アセットID '{asset_id_to_test}' がマニフェストに見つかりません。")

        placeholder1_def = asset_info['placeholders'][0]
        placeholder2_def = asset_info['placeholders'][1]
        
        blueprint = schemas.SlideBlueprint(
            slide_id="test_slide_01",
            slide_title="テスト用組織図",
            asset_id=asset_id_to_test,
            content_map=[
                schemas.PlaceholderContent(placeholder_name=placeholder1_def['name'], content="代表取締役社長"),
                schemas.PlaceholderContent(placeholder_name=placeholder2_def['name'], content="山田 太郎")
            ]
        )

        pptx_io = self.renderer.render_presentation([blueprint])
        self.assertIsNotNone(pptx_io)
        
        try:
            pptx_path = self.output_path / "test_output.pptx"
            with open(pptx_path, "wb") as f:
                f.write(pptx_io.getvalue())
            
            prs = Presentation(pptx_path)
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            
            # --- ★★★ 修正箇所：インデックス(idx)で検証 ---
            placeholder1_idx = placeholder1_def['idx']
            placeholder2_idx = placeholder2_def['idx']

            shape1 = next((s for s in slide.shapes if s.is_placeholder and s.placeholder_format.idx == placeholder1_idx), None)
            shape2 = next((s for s in slide.shapes if s.is_placeholder and s.placeholder_format.idx == placeholder2_idx), None)

            self.assertIsNotNone(shape1, f"プレースホルダー (idx: {placeholder1_idx}) が生成されたスライドに見つかりませんでした。")
            self.assertIsNotNone(shape2, f"プレースホルダー (idx: {placeholder2_idx}) が生成されたスライドに見つかりませんでした。")

            self.assertTrue(shape1.has_text_frame)
            self.assertEqual(shape1.text_frame.text, "代表取締役社長")

            self.assertTrue(shape2.has_text_frame)
            self.assertEqual(shape2.text_frame.text, "山田 太郎")
            # --- 修正ここまで ---

        except Exception as e:
            self.fail(f"生成されたPPTXファイルの検証中にエラーが発生しました: {e}")

if __name__ == '__main__':
    unittest.main()
